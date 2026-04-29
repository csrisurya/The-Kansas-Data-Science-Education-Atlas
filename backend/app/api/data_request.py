"""
API router for data request handling.

Endpoints
---------
POST /data-request      – Submit a new data / dataset request
POST /report/generate   – Generate a PDF or PowerPoint report on demand
"""

import logging
import uuid
import zipfile
from datetime import datetime, timezone
from io import BytesIO
from pathlib import Path
from typing import Optional
from sqlalchemy import text
from app.database import get_db

import pandas as pd
from fastapi import APIRouter, HTTPException, status
from fastapi.responses import FileResponse, StreamingResponse

from app.config import settings
from app.schemas.data_request import (
    DataRequestResponse,
    DataRequestSchema,
    ReportGenerateRequest,
    ReportGenerateResponse,
)
from app.services.email import send_data_request_confirmation

logger = logging.getLogger(__name__)

router = APIRouter()

# Path to raw CSV datasets
DATA_DIR = Path(__file__).resolve().parents[2] / ".." / "data" / "raw"

# Directory where generated ZIPs are saved for download
EXPORTS_DIR = Path(__file__).resolve().parents[2] / ".." / "data" / "exports"
EXPORTS_DIR.resolve().mkdir(parents=True, exist_ok=True)


# Columns that contain county information (used for geographic filtering)
COUNTY_COLUMNS = ["County_Name", "County", "COUNTYNM"]

# ---------------------------------------------------------------------------
# In-memory store (swap for a real DB table when ready)
# ---------------------------------------------------------------------------
_requests_store: list[dict] = []


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _send_confirmation_email(
    name: str,
    email: str,
    request_id: str,
    datasets: list[str] | None = None,
    intended_use: str = "",
    download_url: str = "",
) -> None:
    """
    Best-effort confirmation email via :func:`send_data_request_confirmation`.

    Errors are logged inside the service – this wrapper never raises.
    """
    send_data_request_confirmation(
        recipient_email=email,
        recipient_name=name,
        request_id=request_id,
        datasets=datasets or [],
        intended_use=intended_use,
        download_url=download_url,
    )


# ---------------------------------------------------------------------------
# Helpers – dataset packaging
# ---------------------------------------------------------------------------

# Mapping from dataset key to database table name
DATASET_TABLES = {
    "dataset1": "institutions",
    "dataset2": "county_aggregations",
    "dataset3": "courses",
    "dataset4": "college_locations",
    "dataset5": "digital_infrastructure",
    "dataset6": "socioeconomic",
    "dataset7": "atlas",
}

def _query_table(table_name: str, counties: list[str] | None) -> pd.DataFrame:
    """Query a database table and optionally filter by county."""
    db = next(get_db())
    try:
        with db.bind.connect() as conn:
            df = pd.read_sql(text(f"SELECT * FROM {table_name}"), conn)
        if not counties:
            return df
        for col in COUNTY_COLUMNS:
            if col in df.columns:
                return df[df[col].isin(counties)]
        return df
    finally:
        db.close()

def _build_datasets_zip(
    dataset_keys: list[str],
    data_format: str,
    counties: list[str] | None,
) -> BytesIO:
    """Build a ZIP archive by querying the database for each dataset."""
    zip_buf = BytesIO()
    with zipfile.ZipFile(zip_buf, "w", zipfile.ZIP_DEFLATED) as zf:
        for key in dataset_keys:
            table_name = DATASET_TABLES.get(key)
            if not table_name:
                logger.warning("Unknown dataset key: %s – skipping", key)
                continue
            try:
                df = _query_table(table_name, counties)
                file_bytes, ext = _convert_df(df, data_format)
                arc_name = f"{key}{ext}"
                zf.writestr(arc_name, file_bytes)
                logger.info("Added %s rows from %s to ZIP", len(df), table_name)
            except Exception as e:
                logger.error("Failed to query %s: %s", table_name, e)
                continue
    zip_buf.seek(0)
    return zip_buf


def _convert_df(df: pd.DataFrame, fmt: str) -> tuple[bytes, str]:
    """
    Convert a DataFrame to the requested format.
    Returns (file_bytes, file_extension).
    """
    buf = BytesIO()
    if fmt == "Excel":
        df.to_excel(buf, index=False, engine="openpyxl")
        return buf.getvalue(), ".xlsx"
    elif fmt == "JSON":
        content = df.to_json(orient="records", indent=2)
        return content.encode("utf-8"), ".json"
    else:  # default CSV
        content = df.to_csv(index=False)
        return content.encode("utf-8"), ".csv"



# ---------------------------------------------------------------------------
# POST /data-request
# ---------------------------------------------------------------------------

@router.post(
    "/data-request",
    response_model=DataRequestResponse,
    status_code=status.HTTP_201_CREATED,
    summary="Submit a data request",
    description=(
        "Validates the request, packages the requested datasets into a ZIP archive "
        "in the chosen format, saves it, and sends a confirmation email with a "
        "download link."
    ),
)
async def create_data_request(payload: DataRequestSchema):
    """Accept a data request, build the ZIP, save it, and email the download link."""

    # --- Validation ---
    if payload.request_type == "dataset" and not payload.datasets:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail="At least one dataset must be selected for a dataset request.",
        )
    if payload.request_type == "report" and not payload.report_type:
        raise HTTPException(
            status_code=status.HTTP_422_UNPROCESSABLE_ENTITY,
            detail="Report type is required for a report request.",
        )

    # --- Persist (in-memory for now) ---
    request_id = str(uuid.uuid4())
    record = {
        "request_id": request_id,
        "status": "completed",
        "created_at": datetime.now(timezone.utc).isoformat(),
        **payload.model_dump(),
    }
    _requests_store.append(record)
    logger.info("Data request %s created by %s <%s>", request_id, payload.name, payload.email)

    # --- Build the ZIP of datasets ---
    counties = payload.counties if payload.geographic_scope == "specific" else None
    zip_buf = _build_datasets_zip(payload.datasets, payload.data_format, counties)

    # --- Store ZIP in memory store for download ---
    _requests_store[-1]["zip_data"] = zip_buf.getvalue()
    logger.info("Built dataset ZIP for request %s", request_id)

    # --- Build the download URL ---
    download_url = f"{settings.BACKEND_URL}{settings.API_V1_STR}/data-request/{request_id}/download"

    # --- Send confirmation email with download link ---
    _send_confirmation_email(
        name=payload.name,
        email=payload.email,
        request_id=request_id,
        datasets=payload.datasets,
        intended_use=payload.intended_use or "",
        download_url=download_url,
    )

    return DataRequestResponse(
        request_id=request_id,
        message=(
            f"Your datasets are ready! A download link has been sent to {payload.email}. "
            "Please check your inbox (and spam folder)."
        ),
    )


# ---------------------------------------------------------------------------
# GET /data-request/{request_id}/download
# ---------------------------------------------------------------------------

@router.get(
    "/data-request/{request_id}/download",
    summary="Download a prepared dataset ZIP",
)

async def download_data_request(request_id: str):
    """Serve a previously generated dataset ZIP from memory."""
    record = next((r for r in _requests_store if r["request_id"] == request_id), None)
    if not record or "zip_data" not in record:
        raise HTTPException(
            status_code=status.HTTP_404_NOT_FOUND,
            detail="Download not found. The link may have expired or the request ID is invalid.",
        )
    return StreamingResponse(
        BytesIO(record["zip_data"]),
        media_type="application/zip",
        headers={"Content-Disposition": f'attachment; filename="atlas_data_{request_id[:8]}.zip"'},
    )


# ---------------------------------------------------------------------------
# POST /report/generate
# ---------------------------------------------------------------------------

@router.post(
    "/report/generate",
    summary="Generate a PDF or PowerPoint report",
    description="Builds a report based on selected sections and counties, and returns it as a downloadable file.",
)
async def generate_report(payload: ReportGenerateRequest):
    """Generate and return a report file."""

    if payload.output_format == "PDF":
        return _generate_pdf_report(payload)
    elif payload.output_format == "PowerPoint":
        return _generate_pptx_report(payload)
    else:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Unsupported output format: {payload.output_format}",
        )


def _generate_pdf_report(payload: ReportGenerateRequest) -> StreamingResponse:
    """Create a simple PDF report using fpdf2."""
    try:
        from fpdf import FPDF
    except ImportError:
        raise HTTPException(
            status_code=status.HTTP_501_NOT_IMPLEMENTED,
            detail="PDF generation is not available. Install fpdf2.",
        )

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    # Title
    pdf.set_font("Helvetica", "B", 20)
    pdf.cell(0, 15, "Kansas Data Science Education Atlas", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.ln(5)

    # Report type
    pdf.set_font("Helvetica", "B", 14)
    pdf.cell(0, 10, f"Report: {payload.report_type}", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(3)

    # Counties
    pdf.set_font("Helvetica", "", 11)
    counties_text = "All Counties" if payload.all_counties else ", ".join(payload.counties or [])
    pdf.cell(0, 8, f"Counties: {counties_text}", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(2)

    # Sections
    pdf.set_font("Helvetica", "B", 12)
    pdf.cell(0, 10, "Included Sections:", new_x="LMARGIN", new_y="NEXT")
    pdf.set_font("Helvetica", "", 11)
    for section in payload.sections:
        pdf.cell(0, 8, f"  • {section}", new_x="LMARGIN", new_y="NEXT")
    pdf.ln(5)

    # Placeholder content
    pdf.set_font("Helvetica", "I", 10)
    pdf.multi_cell(0, 7, "Full report content will be populated from the database in a future release.")

    buf = BytesIO(pdf.output())
    buf.seek(0)

    filename = f"atlas_report_{payload.report_type.lower().replace(' ', '_')}.pdf"
    return StreamingResponse(
        buf,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )


def _generate_pptx_report(payload: ReportGenerateRequest) -> StreamingResponse:
    """Create a simple PowerPoint report using python-pptx."""
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
    except ImportError:
        raise HTTPException(
            status_code=status.HTTP_501_NOT_IMPLEMENTED,
            detail="PowerPoint generation is not available. Install python-pptx.",
        )

    prs = Presentation()

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Kansas Data Science Education Atlas"
    slide.placeholders[1].text = f"Report: {payload.report_type}"

    # Content slide
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = "Report Summary"

    body = slide.placeholders[1]
    tf = body.text_frame
    tf.text = f"Counties: {'All Counties' if payload.all_counties else ', '.join(payload.counties or [])}"

    for section in payload.sections:
        p = tf.add_paragraph()
        p.text = f"• {section}"

    p = tf.add_paragraph()
    p.text = ""
    p = tf.add_paragraph()
    p.text = "Full report content will be populated from the database in a future release."

    buf = BytesIO()
    prs.save(buf)
    buf.seek(0)

    filename = f"atlas_report_{payload.report_type.lower().replace(' ', '_')}.pptx"
    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="{filename}"'},
    )
